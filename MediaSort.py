"""
MediaSort — sorts images, videos, documents, and audio into organised sub-folders.
             Includes an AI Rename tab powered by a local GGUF model via llama-server.

┌─────────────────────────────────────────────────────────────────┐
│  IMAGE categories  (checked in priority order)                  │
│   1. Wallpapers      → filename contains "wallhaven"            │
│   2. Screenshots     → filename contains "screenshot"           │
│   3. WebP Images     → extension is .webp                       │
│   4. Square Images   → width/height ratio within 5 % of 1:1    │
│   5. Other Images    → everything else                          │
├─────────────────────────────────────────────────────────────────┤
│  VIDEO categories  (checked in priority order)                  │
│   1. Screen Recordings → filename starts with "screen_recording"│
│   2. 16:9 Landscape    → aspect ratio ≈ 16:9                    │
│   3. 9:16 Portrait     → aspect ratio ≈ 9:16                    │
│   4. Other             → anything else                          │
├─────────────────────────────────────────────────────────────────┤
│  DOCUMENT categories                                            │
│   1. Word Documents  → .doc .docx .odt .rtf                    │
│   2. Spreadsheets    → .xls .xlsx .ods .csv .tsv               │
│   3. Presentations   → .ppt .pptx .odp                         │
│   4. PDFs            → .pdf                                     │
│   5. Text & Markdown → .txt .md .rst .log                      │
│   6. Code & Scripts  → .py .js .ts .html .css .json .xml ...   │
│   7. Archives        → .zip .rar .7z .tar .gz .bz2 .xz         │
│   8. Other Docs      → everything else                          │
├─────────────────────────────────────────────────────────────────┤
│  AUDIO categories                                               │
│   1. Lossless        → .flac .wav .aiff .aif .alac .ape .wv    │
│   2. Compressed      → .mp3 .aac .ogg .opus .m4a .wma .ac3 ... │
│   3. Playlists       → .m3u .m3u8 .pls .xspf .wpl              │
│   4. Other Audio     → everything else                          │
├─────────────────────────────────────────────────────────────────┤
│  AI RENAME  (supported: .txt .pdf .docx .pptx)                 │
│   • Reads each file page-by-page with the best library          │
│   • Sends page summaries to a local llama-server (GGUF/Vulkan)  │
│   • Model suggests a short filename in the document's language  │
│   • User reviews suggestions and applies renames with one click │
│   • llama-server endpoint configurable in the UI                │
└─────────────────────────────────────────────────────────────────┘

Dependencies:
  tkinter            — GUI (stdlib, bundled with Python on Windows)
  threading          — background scan (stdlib)
  os, pathlib, shutil, math, subprocess, json, urllib  (stdlib)
  Pillow             — image dimensions        (pip install Pillow)
  ffprobe (on PATH)  — video dimensions        (or opencv-python as fallback)
  opencv-python      — optional fallback       (pip install opencv-python)
  pypdf              — PDF text extraction     (pip install pypdf)
  python-docx        — Word text extraction    (pip install python-docx)
  python-pptx        — PowerPoint extraction   (pip install python-pptx)
  llama-server       — local GGUF inference    (llama.cpp Vulkan build on PATH)
"""

import os
import json
import math
import shutil
import subprocess
import threading
import urllib.request
import urllib.error
from pathlib import Path
import tkinter as tk
from tkinter import ttk, filedialog, messagebox

try:
    from PIL import Image as _PILImage
    _PIL_AVAILABLE = True
except ImportError:
    _PIL_AVAILABLE = False

try:
    import pypdf as _pypdf
    _PYPDF_AVAILABLE = True
except ImportError:
    _PYPDF_AVAILABLE = False

try:
    from docx import Document as _DocxDocument
    _DOCX_AVAILABLE = True
except ImportError:
    _DOCX_AVAILABLE = False

try:
    from pptx import Presentation as _PptxPresentation
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False


# ══════════════════════════════════════════════════════════════════
# Shared palette
# ══════════════════════════════════════════════════════════════════

PALETTE = {
    "bg":        "#0f0f13",
    "panel":     "#17171f",
    "card":      "#1e1e2a",
    "border":    "#2a2a3a",
    "accent":    "#6c63ff",
    "accent2":   "#ff6584",
    "accent3":   "#43e8a0",
    "accent4":   "#ffd166",
    "text":      "#e8e8f0",
    "subtext":   "#888898",
    "hover":     "#252535",
    "sel":       "#2d2d4a",
    # image category colours
    "wall":      "#b39ddb",   # soft purple  — wallpapers
    "shot":      "#80cbc4",   # teal         — screenshots
    "webp":      "#ffb74d",   # amber        — webp
    "square":    "#4dd0e1",   # cyan         — square images
    "img_other": "#ef9a9a",   # rose         — other images
    # video category colours
    "landscape": "#4fc3f7",   # sky blue     — 16:9
    "portrait":  "#f48fb1",   # pink         — 9:16
    "screen":    "#a5d6a7",   # green        — screen recordings
    "vid_other": "#ffcc80",   # peach        — other videos
    # document category colours
    "word":      "#90caf9",   # light blue   — word docs
    "sheet":     "#a5d6a7",   # light green  — spreadsheets
    "slides":    "#ffcc02",   # yellow       — presentations
    "pdf":       "#ef9a9a",   # rose         — PDFs
    "textmd":    "#ce93d8",   # lilac        — text & markdown
    "code":      "#80deea",   # cyan         — code & scripts
    "archive":   "#ffab91",   # orange       — archives
    "doc_other": "#b0bec5",   # blue-grey    — other docs
    # audio category colours
    "lossless":  "#e6ee9c",   # lime         — lossless audio
    "lossy":     "#80cbc4",   # teal         — compressed audio
    "playlist":  "#b39ddb",   # purple       — playlists
    "aud_other": "#f48fb1",   # pink         — other audio
    # ai rename colours
    "ai":        "#c792ea",   # violet       — AI rename accent
    "ai_ok":     "#43e8a0",   # green        — suggestion accepted
    "ai_skip":   "#888898",   # grey         — skipped
    "ai_err":    "#ff6584",   # red          — error
    "ai_pend":   "#ffd166",   # amber        — pending
}


# ══════════════════════════════════════════════════════════════════
# IMAGE constants & logic
# ══════════════════════════════════════════════════════════════════

IMAGE_EXTENSIONS = {
    ".jpg", ".jpeg", ".png", ".webp", ".gif", ".bmp",
    ".tiff", ".tif", ".heic", ".heif", ".avif", ".jfif",
    ".ico", ".svg", ".raw", ".cr2", ".nef", ".arw",
    ".dng", ".orf", ".rw2", ".psd", ".xcf",
}

ICАТ_WALLPAPER  = "Wallpapers"
ICАТ_SCREENSHOT = "Screenshots"
ICАТ_WEBP       = "WebP Images"
ICАТ_SQUARE     = "Square Images"
ICАТ_OTHER      = "Other Images"

IMAGE_CATEGORIES = [ICАТ_WALLPAPER, ICАТ_SCREENSHOT, ICАТ_WEBP, ICАТ_SQUARE, ICАТ_OTHER]

# 5 % tolerance: shorter side must be >= 95 % of longer side
SQUARE_TOLERANCE = 0.05

IMAGE_CAT_COLORS = {
    ICАТ_WALLPAPER:  PALETTE["wall"],
    ICАТ_SCREENSHOT: PALETTE["shot"],
    ICАТ_WEBP:       PALETTE["webp"],
    ICАТ_SQUARE:     PALETTE["square"],
    ICАТ_OTHER:      PALETTE["img_other"],
}

IMAGE_CAT_ICONS = {
    ICАТ_WALLPAPER:  "🖼",
    ICАТ_SCREENSHOT: "📸",
    ICАТ_WEBP:       "⚡",
    ICАТ_SQUARE:     "⬛",
    ICАТ_OTHER:      "📁",
}


def _get_image_dimensions(filepath: Path):
    """Return (width, height) via Pillow, or None if unavailable/unreadable."""
    if not _PIL_AVAILABLE:
        return None
    try:
        with _PILImage.open(filepath) as img:
            return img.size  # (width, height)
    except Exception:
        return None


def classify_image(filepath: Path) -> str:
    """Classify one image file. Priority order documented in module docstring."""
    name_lower = filepath.stem.lower()
    ext_lower  = filepath.suffix.lower()

    if "wallhaven" in name_lower:
        return ICАТ_WALLPAPER
    if "screenshot" in name_lower:
        return ICАТ_SCREENSHOT
    if ext_lower == ".webp":
        return ICАТ_WEBP

    dims = _get_image_dimensions(filepath)
    if dims is not None:
        w, h = dims
        if w > 0 and h > 0:
            ratio = min(w, h) / max(w, h)
            if ratio >= (1.0 - SQUARE_TOLERANCE):
                return ICАТ_SQUARE

    return ICАТ_OTHER


# ══════════════════════════════════════════════════════════════════
# VIDEO constants & logic
# ══════════════════════════════════════════════════════════════════

VIDEO_EXTENSIONS = {
    ".mp4", ".mkv", ".avi", ".mov", ".wmv", ".flv", ".webm",
    ".m4v", ".ts", ".mts", ".m2ts", ".vob", ".3gp", ".3g2",
    ".mpg", ".mpeg", ".f4v", ".rm", ".rmvb", ".divx", ".ogv",
}

VCAT_SCREEN_REC = "Screen Recordings"
VCAT_LANDSCAPE  = "16:9 Landscape"
VCAT_PORTRAIT   = "9:16 Portrait"
VCAT_OTHER      = "Other"

VIDEO_CATEGORIES = [VCAT_SCREEN_REC, VCAT_LANDSCAPE, VCAT_PORTRAIT, VCAT_OTHER]

RATIO_TOLERANCE = 0.05   # ±5 % when comparing aspect ratios

VIDEO_CAT_COLORS = {
    VCAT_SCREEN_REC: PALETTE["screen"],
    VCAT_LANDSCAPE:  PALETTE["landscape"],
    VCAT_PORTRAIT:   PALETTE["portrait"],
    VCAT_OTHER:      PALETTE["vid_other"],
}

VIDEO_CAT_ICONS = {
    VCAT_SCREEN_REC: "🎬",
    VCAT_LANDSCAPE:  "🖥",
    VCAT_PORTRAIT:   "📱",
    VCAT_OTHER:      "📂",
}


def _get_dims_ffprobe(path: str):
    """Return (width, height) via ffprobe, or None on failure."""
    cmd = ["ffprobe", "-v", "quiet", "-print_format", "json", "-show_streams", path]
    try:
        result = subprocess.run(cmd, capture_output=True, timeout=10)
        data = json.loads(result.stdout)
        for stream in data.get("streams", []):
            w = stream.get("width")
            h = stream.get("height")
            if w and h:
                return int(w), int(h)
    except Exception:
        pass
    return None


def _get_dims_cv2(path: str):
    """Return (width, height) via OpenCV, or None on failure."""
    try:
        import cv2
        cap = cv2.VideoCapture(path)
        if cap.isOpened():
            w = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
            h = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
            cap.release()
            if w and h:
                return w, h
    except Exception:
        pass
    return None


def get_video_dimensions(path: str):
    """Try ffprobe first, then cv2 as fallback. Returns (w, h) or None."""
    dims = _get_dims_ffprobe(path)
    if dims is None:
        dims = _get_dims_cv2(path)
    return dims


def classify_video(filepath: Path) -> str:
    """Classify one video file. Priority order documented in module docstring."""
    name_lower = filepath.name.lower()

    if name_lower.startswith("screen_recording"):
        return VCAT_SCREEN_REC

    dims = get_video_dimensions(str(filepath))
    if dims is None:
        return VCAT_OTHER

    w, h = dims
    if w == 0 or h == 0:
        return VCAT_OTHER

    ratio      = w / h
    target_169 = 16 / 9
    target_916 = 9  / 16

    if abs(ratio - target_169) / target_169 <= RATIO_TOLERANCE:
        return VCAT_LANDSCAPE
    if abs(ratio - target_916) / target_916 <= RATIO_TOLERANCE:
        return VCAT_PORTRAIT

    return VCAT_OTHER


# ══════════════════════════════════════════════════════════════════
# DOCUMENT constants & logic
# ══════════════════════════════════════════════════════════════════

DCAT_WORD     = "Word Documents"
DCAT_SHEET    = "Spreadsheets"
DCAT_SLIDES   = "Presentations"
DCAT_PDF      = "PDFs"
DCAT_TEXTMD   = "Text & Markdown"
DCAT_CODE     = "Code & Scripts"
DCAT_ARCHIVE  = "Archives"
DCAT_OTHER    = "Other Documents"

DOCUMENT_CATEGORIES = [
    DCAT_WORD, DCAT_SHEET, DCAT_SLIDES, DCAT_PDF,
    DCAT_TEXTMD, DCAT_CODE, DCAT_ARCHIVE, DCAT_OTHER,
]

_DOC_EXT_MAP: dict[str, str] = {}

def _doc_exts(*exts, cat):
    for e in exts:
        _DOC_EXT_MAP[e] = cat

_doc_exts(".doc", ".docx", ".odt", ".rtf", ".dot", ".dotx", ".docm",        cat=DCAT_WORD)
_doc_exts(".xls", ".xlsx", ".ods", ".csv", ".tsv", ".xlsm", ".xlsb",        cat=DCAT_SHEET)
_doc_exts(".ppt", ".pptx", ".odp", ".pot", ".potx", ".pptm",                cat=DCAT_SLIDES)
_doc_exts(".pdf",                                                             cat=DCAT_PDF)
_doc_exts(".txt", ".md", ".rst", ".log", ".nfo", ".readme", ".text",         cat=DCAT_TEXTMD)
_doc_exts(
    ".py", ".js", ".ts", ".jsx", ".tsx", ".html", ".htm", ".css", ".scss",
    ".json", ".xml", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".env",
    ".sql", ".sh", ".bat", ".cmd", ".ps1", ".rb", ".php", ".java", ".c",
    ".cpp", ".cs", ".go", ".rs", ".kt", ".swift", ".r", ".lua", ".pl",
                                                                              cat=DCAT_CODE)
_doc_exts(".zip", ".rar", ".7z", ".tar", ".gz", ".bz2", ".xz",
          ".tgz", ".tbz2", ".lz", ".lzma", ".zst", ".cab", ".iso",           cat=DCAT_ARCHIVE)

DOCUMENT_EXTENSIONS = set(_DOC_EXT_MAP.keys()) | {DCAT_OTHER}

def classify_document(filepath: Path) -> str:
    return _DOC_EXT_MAP.get(filepath.suffix.lower(), DCAT_OTHER)

DOCUMENT_CAT_COLORS = {
    DCAT_WORD:    PALETTE["word"],
    DCAT_SHEET:   PALETTE["sheet"],
    DCAT_SLIDES:  PALETTE["slides"],
    DCAT_PDF:     PALETTE["pdf"],
    DCAT_TEXTMD:  PALETTE["textmd"],
    DCAT_CODE:    PALETTE["code"],
    DCAT_ARCHIVE: PALETTE["archive"],
    DCAT_OTHER:   PALETTE["doc_other"],
}

DOCUMENT_CAT_ICONS = {
    DCAT_WORD:    "📝",
    DCAT_SHEET:   "📊",
    DCAT_SLIDES:  "📑",
    DCAT_PDF:     "📄",
    DCAT_TEXTMD:  "📃",
    DCAT_CODE:    "💻",
    DCAT_ARCHIVE: "🗜",
    DCAT_OTHER:   "📁",
}


# ══════════════════════════════════════════════════════════════════
# AUDIO constants & logic
# ══════════════════════════════════════════════════════════════════

ACAT_LOSSLESS = "Lossless"
ACAT_LOSSY    = "Compressed"
ACAT_PLAYLIST = "Playlists"
ACAT_OTHER    = "Other Audio"

AUDIO_CATEGORIES = [ACAT_LOSSLESS, ACAT_LOSSY, ACAT_PLAYLIST, ACAT_OTHER]

_AUDIO_EXT_MAP: dict[str, str] = {}

def _aud_exts(*exts, cat):
    for e in exts:
        _AUDIO_EXT_MAP[e] = cat

_aud_exts(".flac", ".wav", ".aiff", ".aif", ".alac", ".ape",
          ".wv", ".tta", ".dsd", ".dsf", ".dff", ".pcm",       cat=ACAT_LOSSLESS)
_aud_exts(".mp3", ".aac", ".ogg", ".opus", ".m4a", ".wma",
          ".ac3", ".amr", ".mp2", ".mka", ".ra", ".mid",
          ".midi", ".spx", ".voc", ".au",                       cat=ACAT_LOSSY)
_aud_exts(".m3u", ".m3u8", ".pls", ".xspf", ".wpl", ".asx",    cat=ACAT_PLAYLIST)

AUDIO_EXTENSIONS = set(_AUDIO_EXT_MAP.keys())

def classify_audio(filepath: Path) -> str:
    return _AUDIO_EXT_MAP.get(filepath.suffix.lower(), ACAT_OTHER)

AUDIO_CAT_COLORS = {
    ACAT_LOSSLESS: PALETTE["lossless"],
    ACAT_LOSSY:    PALETTE["lossy"],
    ACAT_PLAYLIST: PALETTE["playlist"],
    ACAT_OTHER:    PALETTE["aud_other"],
}

AUDIO_CAT_ICONS = {
    ACAT_LOSSLESS: "🎵",
    ACAT_LOSSY:    "🎧",
    ACAT_PLAYLIST: "📋",
    ACAT_OTHER:    "🔊",
}


# ══════════════════════════════════════════════════════════════════
# Shared helpers
# ══════════════════════════════════════════════════════════════════

def _human_size(nbytes: int) -> str:
    for unit in ("B", "KB", "MB", "GB"):
        if nbytes < 1024:
            return f"{nbytes:.0f} {unit}" if unit == "B" else f"{nbytes:.1f} {unit}"
        nbytes /= 1024
    return f"{nbytes:.1f} TB"


def _safe_folder_name(name: str) -> str:
    """Strip characters Windows forbids in directory names."""
    for ch in r'\/:*?"<>|':
        name = name.replace(ch, "-")
    return name.strip()


def _make_btn(parent, text, cmd, color):
    """Flat styled button with hover effect, shared by both sorter panels."""
    b = tk.Button(
        parent, text=text, command=cmd,
        font=("Segoe UI", 9, "bold"),
        fg=color, bg=PALETTE["card"],
        activeforeground=PALETTE["text"], activebackground=PALETTE["hover"],
        relief="flat", bd=0, cursor="hand2", padx=12, pady=5,
    )
    b.bind("<Enter>", lambda e: b.configure(bg=PALETTE["hover"]))
    b.bind("<Leave>", lambda e: b.configure(bg=PALETTE["card"]))
    return b


def _apply_shared_styles(widget):
    """Apply dark ttk theme. Call once after the first Tk window is created."""
    s = ttk.Style(widget)
    s.theme_use("clam")
    s.configure("Treeview",
                background=PALETTE["panel"], foreground=PALETTE["text"],
                fieldbackground=PALETTE["panel"], borderwidth=0,
                font=("Segoe UI", 9), rowheight=24)
    s.configure("Treeview.Heading",
                background=PALETTE["card"], foreground=PALETTE["subtext"],
                borderwidth=0, font=("Segoe UI", 9, "bold"))
    s.map("Treeview",
          background=[("selected", PALETTE["sel"])],
          foreground=[("selected", PALETTE["text"])])
    s.configure("TNotebook", background=PALETTE["bg"], borderwidth=0)
    s.configure("TNotebook.Tab",
                background=PALETTE["card"], foreground=PALETTE["subtext"],
                padding=[14, 6], font=("Segoe UI", 9, "bold"), borderwidth=0)
    s.map("TNotebook.Tab",
          background=[("selected", PALETTE["panel"])],
          foreground=[("selected", PALETTE["text"])])
    s.configure("TProgressbar",
                troughcolor=PALETTE["card"], background=PALETTE["accent"],
                borderwidth=0, thickness=6)
    s.configure("TScrollbar",
                background=PALETTE["border"], troughcolor=PALETTE["panel"],
                borderwidth=0, arrowsize=0)


# ══════════════════════════════════════════════════════════════════
# Base sorter panel (shared scaffold for both Image and Video tabs)
# ══════════════════════════════════════════════════════════════════

class _SorterPanel(tk.Frame):
    """
    Abstract base for a sorter tab. Subclasses must set:
        self._categories   : list[str]
        self._cat_colors   : dict[str, str]
        self._cat_icons    : dict[str, str]
        self._extensions   : set[str]
    and implement:
        _classify(filepath)  → str   category name
        _scan_row(fp, root)  → tuple  values for the treeview row
        _tree_columns()      → list[tuple(id, heading, width, anchor)]
    """

    def __init__(self, parent):
        super().__init__(parent, bg=PALETTE["bg"])
        # State
        self._folder    = tk.StringVar()
        self._recursive = tk.BooleanVar(value=True)
        self._auto_move = tk.BooleanVar(value=False)
        self._status    = tk.StringVar(value="Choose a folder to begin.")
        self._progress  = tk.DoubleVar(value=0.0)
        self._results: dict[str, list[Path]] = {}
        self._scanning  = False
        self._moving    = False
        self._stop_flag = threading.Event()

    def _init_results(self):
        self._results = {c: [] for c in self._categories}

    # ── Scaffold builder ──────────────────────

    def build(self, subtitle: str):
        """Call from subclass __init__ after setting constants."""
        self._init_results()
        self._build_toolbar(subtitle)
        self._build_options()
        self._build_progress()
        self._build_cards()
        self._build_notebook()
        self._build_action_bar()

    def _build_toolbar(self, subtitle: str):
        picker = tk.Frame(self, bg=PALETTE["card"])
        picker.pack(fill="x", padx=20, pady=(14, 0), ipady=10)

        tk.Label(picker, text="Folder:", font=("Segoe UI", 10, "bold"),
                 fg=PALETTE["subtext"], bg=PALETTE["card"]).pack(side="left", padx=(14, 6))

        tk.Entry(
            picker, textvariable=self._folder,
            font=("Segoe UI", 10), fg=PALETTE["text"],
            bg=PALETTE["panel"], bd=0, insertbackground=PALETTE["text"],
            relief="flat", width=58,
        ).pack(side="left", ipady=5, padx=(0, 8), fill="x", expand=True)

        _make_btn(picker, "Browse",   self._browse,      PALETTE["accent"]).pack(side="left", padx=(0, 6))

        self._scan_btn = _make_btn(picker, "▶  Scan", self._start_scan, PALETTE["accent3"])
        self._scan_btn.pack(side="left", padx=(0, 6))

        self._stop_btn = _make_btn(picker, "■  Stop", self._stop_scan, PALETTE["accent2"])
        self._stop_btn.pack(side="left", padx=(0, 14))
        self._stop_btn.configure(state="disabled")

    def _build_options(self):
        opts = tk.Frame(self, bg=PALETTE["bg"])
        opts.pack(fill="x", padx=20, pady=(8, 0))

        tk.Checkbutton(
            opts, text="Scan sub-folders recursively",
            variable=self._recursive,
            font=("Segoe UI", 9), fg=PALETTE["subtext"],
            bg=PALETTE["bg"], activebackground=PALETTE["bg"],
            activeforeground=PALETTE["text"], selectcolor=PALETTE["card"],
        ).pack(side="left")

        tk.Label(opts, text=" │ ", font=("Segoe UI", 9),
                 fg=PALETTE["border"], bg=PALETTE["bg"]).pack(side="left")

        tk.Checkbutton(
            opts, text="Auto-move files into sub-folders after scan",
            variable=self._auto_move,
            font=("Segoe UI", 9, "bold"), fg=PALETTE["accent2"],
            bg=PALETTE["bg"], activebackground=PALETTE["bg"],
            activeforeground=PALETTE["text"], selectcolor=PALETTE["card"],
        ).pack(side="left")

    def _build_progress(self):
        prog_frame = tk.Frame(self, bg=PALETTE["bg"])
        prog_frame.pack(fill="x", padx=20, pady=(10, 0))

        ttk.Progressbar(
            prog_frame, variable=self._progress,
            maximum=100, mode="determinate",
        ).pack(side="left", fill="x", expand=True)

        tk.Label(
            prog_frame, textvariable=self._status,
            font=("Segoe UI", 9), fg=PALETTE["subtext"],
            bg=PALETTE["bg"], anchor="w",
        ).pack(side="left", padx=(10, 0))

    def _build_cards(self):
        cards = tk.Frame(self, bg=PALETTE["bg"])
        cards.pack(fill="x", padx=20, pady=(14, 0))
        self._count_labels: dict[str, tk.Label] = {}

        for cat in self._categories:
            card = tk.Frame(cards, bg=PALETTE["card"])
            card.pack(side="left", expand=True, fill="x", padx=(0, 8), ipady=10)

            color = self._cat_colors[cat]
            tk.Frame(card, bg=color, width=4).pack(side="left", fill="y")

            inner = tk.Frame(card, bg=PALETTE["card"])
            inner.pack(side="left", padx=12, pady=6)

            title_row = tk.Frame(inner, bg=PALETTE["card"])
            title_row.pack(anchor="w")

            tk.Label(title_row, text=self._cat_icons[cat] + " ",
                     font=("Segoe UI", 11), fg=color, bg=PALETTE["card"]).pack(side="left")
            tk.Label(title_row, text=cat,
                     font=("Segoe UI", 9, "bold"), fg=color, bg=PALETTE["card"]).pack(side="left")

            cnt = tk.Label(inner, text="0 files",
                           font=("Segoe UI", 18, "bold"),
                           fg=PALETTE["text"], bg=PALETTE["card"])
            cnt.pack(anchor="w")
            self._count_labels[cat] = cnt

    def _build_notebook(self):
        nb_frame = tk.Frame(self, bg=PALETTE["bg"])
        nb_frame.pack(fill="both", expand=True, padx=20, pady=(14, 0))

        self._nb = ttk.Notebook(nb_frame)
        self._nb.pack(fill="both", expand=True)

        self._trees: dict[str, ttk.Treeview] = {}
        cols_spec = self._tree_columns()

        for cat in self._categories:
            tab = tk.Frame(self._nb, bg=PALETTE["panel"])
            self._nb.add(tab, text=f"  {self._cat_icons[cat]}  {cat}  ")

            col_ids = [c[0] for c in cols_spec]
            tree = ttk.Treeview(tab, columns=col_ids, show="headings", selectmode="browse")

            for col_id, heading, width, anchor in cols_spec:
                tree.heading(col_id, text=heading)
                tree.column(col_id, width=width, anchor=anchor)

            vsb = ttk.Scrollbar(tab, orient="vertical",   command=tree.yview)
            hsb = ttk.Scrollbar(tab, orient="horizontal", command=tree.xview)
            tree.configure(yscrollcommand=vsb.set, xscrollcommand=hsb.set)

            vsb.grid(row=0, column=1, sticky="ns")
            hsb.grid(row=1, column=0, sticky="ew")
            tree.grid(row=0, column=0, sticky="nsew")
            tab.rowconfigure(0, weight=1)
            tab.columnconfigure(0, weight=1)

            self._trees[cat] = tree

    def _build_action_bar(self):
        action_bar = tk.Frame(self, bg=PALETTE["panel"])
        action_bar.pack(fill="x", padx=20, pady=(10, 16))

        tk.Label(action_bar, text="Actions:",
                 font=("Segoe UI", 9, "bold"),
                 fg=PALETTE["subtext"], bg=PALETTE["panel"]).pack(side="left", padx=(14, 10))

        self._move_btn = _make_btn(
            action_bar, "Move files into sub-folders",
            lambda: self._organize(move=True), PALETTE["accent2"])
        self._move_btn.pack(side="left", padx=(0, 6))
        self._move_btn.configure(state="disabled")

        self._copy_btn = _make_btn(
            action_bar, "Copy files into sub-folders",
            lambda: self._organize(move=False), PALETTE["accent"])
        self._copy_btn.pack(side="left")
        self._copy_btn.configure(state="disabled")

        tk.Label(action_bar,
                 text="Sub-folders are created inside the scanned folder.",
                 font=("Segoe UI", 8), fg=PALETTE["subtext"], bg=PALETTE["panel"],
                 ).pack(side="right", padx=14)

    # ── Browse / Scan ─────────────────────────

    def _browse(self):
        folder = filedialog.askdirectory(title="Select folder to scan")
        if folder:
            self._folder.set(folder)

    def _start_scan(self):
        folder = self._folder.get().strip()
        if not folder or not os.path.isdir(folder):
            messagebox.showerror("No folder", "Please select a valid folder first.")
            return
        if self._scanning:
            return

        self._init_results()
        for cat in self._categories:
            self._trees[cat].delete(*self._trees[cat].get_children())
            self._count_labels[cat].configure(text="0 files")

        self._stop_flag.clear()
        self._scanning = True
        self._scan_btn.configure(state="disabled")
        self._stop_btn.configure(state="normal")
        self._move_btn.configure(state="disabled")
        self._copy_btn.configure(state="disabled")
        self._progress.set(0)

        threading.Thread(target=self._scan_worker, args=(folder,), daemon=True).start()

    def _stop_scan(self):
        self._stop_flag.set()

    def _scan_worker(self, folder: str):
        try:
            file_type = self._file_type_label()
            self._set_status(f"Collecting {file_type} files…")
            root = Path(folder)

            if self._recursive.get():
                all_files = [p for p in root.rglob("*")
                             if p.is_file() and p.suffix.lower() in self._extensions]
            else:
                all_files = [p for p in root.iterdir()
                             if p.is_file() and p.suffix.lower() in self._extensions]

            total = len(all_files)
            if total == 0:
                self._set_status(f"No {file_type} files found.")
                self.after(0, self._finish_scan)
                return

            self._set_status(f"Found {total} {file_type}(s). Classifying…")

            for i, fp in enumerate(all_files):
                if self._stop_flag.is_set():
                    self.after(0, self._set_status, f"Stopped at {i}/{total}.")
                    break

                self.after(0, self._progress.set, (i / total) * 100)
                self.after(0, self._set_status, f"[{i+1}/{total}]  {fp.name}")

                try:
                    cat     = self._classify(fp)
                    row_vals = self._scan_row(fp, root)
                except Exception:
                    cat      = self._categories[-1]   # fallback: last cat = "Other"
                    row_vals = self._error_row(fp, root)

                self._results[cat].append(fp)
                self.after(0, self._add_row, cat, row_vals)
                self.after(0, self._refresh_count, cat)

            self.after(0, self._progress.set, 100)
            found = sum(len(v) for v in self._results.values())
            self.after(0, self._set_status, f"Done — {found} {file_type}(s) classified.")

        except Exception as exc:
            self.after(0, self._set_status, f"Error: {exc}")
        finally:
            self.after(0, self._finish_scan)

    def _add_row(self, cat: str, vals: tuple):
        tree = self._trees[cat]
        iid = tree.insert("", "end", values=vals)
        tag = "even" if len(tree.get_children()) % 2 == 0 else "odd"
        tree.item(iid, tags=(tag,))
        tree.tag_configure("odd",  background=PALETTE["panel"])
        tree.tag_configure("even", background=PALETTE["card"])

    def _refresh_count(self, cat: str):
        n = len(self._results[cat])
        self._count_labels[cat].configure(text=f"{n} file{'s' if n != 1 else ''}")

    def _finish_scan(self):
        self._scanning = False
        self._scan_btn.configure(state="normal")
        self._stop_btn.configure(state="disabled")
        has = any(self._results[c] for c in self._categories)
        st  = "normal" if has else "disabled"
        self._move_btn.configure(state=st)
        self._copy_btn.configure(state=st)
        if has and self._auto_move.get() and not self._stop_flag.is_set():
            self.after(200, lambda: self._organize(move=True, silent=True))

    # ── Organize ──────────────────────────────

    def _organize(self, move: bool, silent: bool = False):
        folder = self._folder.get().strip()
        if not folder or self._moving:
            return

        total = sum(len(v) for v in self._results.values())
        if total == 0:
            if not silent:
                messagebox.showinfo("Nothing to do", "No files to organise.")
            return

        verb = "move" if move else "copy"
        if not silent:
            preview = "\n".join(
                f"  {self._cat_icons[c]}  {c}/  ({len(self._results[c])} files)"
                for c in self._categories if self._results[c]
            )
            if not messagebox.askyesno(
                "Confirm",
                f"This will {verb} {total} file(s) into sub-folders inside:\n{folder}\n\n"
                f"Destination folders:\n{preview}\n\nProceed?"
            ):
                return

        self._moving = True
        self._move_btn.configure(state="disabled")
        self._copy_btn.configure(state="disabled")
        self._scan_btn.configure(state="disabled")
        self._progress.set(0)

        threading.Thread(target=self._move_worker, args=(folder, move), daemon=True).start()

    def _move_worker(self, folder: str, move: bool):
        root   = Path(folder)
        verb   = "Moving" if move else "Copying"
        errors: list[str] = []
        done   = 0

        work: list[tuple[Path, Path]] = []
        for cat, files in self._results.items():
            if not files:
                continue
            dest_dir = root / _safe_folder_name(cat)
            try:
                dest_dir.mkdir(exist_ok=True)
            except Exception as e:
                errors.append(f"mkdir {cat}: {e}")
                continue

            for fp in files:
                dest = dest_dir / fp.name
                if dest.exists() and dest.resolve() != fp.resolve():
                    stem, suffix = fp.stem, fp.suffix
                    n = 1
                    while dest.exists():
                        dest = dest_dir / f"{stem}_{n}{suffix}"
                        n += 1
                work.append((fp, dest))

        actual = len(work)
        for i, (src, dst) in enumerate(work):
            self.after(0, self._progress.set, (i / actual * 100) if actual else 100)
            self.after(0, self._set_status,
                       f"{verb} [{i+1}/{actual}]  {src.name}  →  {dst.parent.name}/")
            try:
                if src.resolve() == dst.resolve():
                    done += 1
                    continue
                if move:
                    shutil.move(str(src), str(dst))
                else:
                    shutil.copy2(str(src), str(dst))
                done += 1
            except Exception as e:
                errors.append(f"{src.name}: {e}")

        self.after(0, self._progress.set, 100)
        action = "moved" if move else "copied"
        lines  = [f"✓  {done} file(s) {action} successfully.\n"]
        for cat in self._categories:
            n = len(self._results[cat])
            if n:
                lines.append(f"  {self._cat_icons[cat]}  {_safe_folder_name(cat)}/  ← {n} file(s)")
        if errors:
            lines.append(f"\n✗  {len(errors)} error(s):")
            lines.extend(f"   • {e}" for e in errors[:15])
            if len(errors) > 15:
                lines.append(f"   … and {len(errors)-15} more.")

        self.after(0, self._set_status, f"Done — {done}/{actual} files {action}.")
        self.after(0, messagebox.showinfo, "Done", "\n".join(lines))
        self.after(0, self._finish_move)

    def _finish_move(self):
        self._moving = False
        self._scan_btn.configure(state="normal")
        has = any(self._results[c] for c in self._categories)
        st  = "normal" if has else "disabled"
        self._move_btn.configure(state=st)
        self._copy_btn.configure(state=st)

    # ── Helpers ───────────────────────────────

    def _set_status(self, msg: str):
        self._status.set(msg)

    def _file_type_label(self) -> str:
        return "file"

    def _error_row(self, fp: Path, root: Path) -> tuple:
        rel = str(fp.parent.relative_to(root)) if fp.parent != root else "."
        return (fp.name, fp.suffix.lower(), rel, "—")

    # ── Abstract interface ────────────────────

    def _classify(self, filepath: Path) -> str:
        raise NotImplementedError

    def _scan_row(self, fp: Path, root: Path) -> tuple:
        raise NotImplementedError

    def _tree_columns(self) -> list:
        raise NotImplementedError


# ══════════════════════════════════════════════════════════════════
# Image Sorter panel
# ══════════════════════════════════════════════════════════════════

class ImageSorterPanel(_SorterPanel):

    def __init__(self, parent):
        super().__init__(parent)
        self._categories  = IMAGE_CATEGORIES
        self._cat_colors  = IMAGE_CAT_COLORS
        self._cat_icons   = IMAGE_CAT_ICONS
        self._extensions  = IMAGE_EXTENSIONS
        self.build(subtitle="wallpapers · screenshots · webp · square · everything else")

    def _file_type_label(self):
        return "image"

    def _classify(self, filepath: Path) -> str:
        return classify_image(filepath)

    def _tree_columns(self):
        return [
            ("name",   "File Name", 300, "w"),
            ("ext",    "Ext",        60, "center"),
            ("folder", "Folder",    340, "w"),
            ("size",   "Size",       80, "e"),
        ]

    def _scan_row(self, fp: Path, root: Path) -> tuple:
        size_str = _human_size(fp.stat().st_size)
        rel      = str(fp.parent.relative_to(root)) if fp.parent != root else "."
        return (fp.name, fp.suffix.lower(), rel, size_str)

    def _error_row(self, fp: Path, root: Path) -> tuple:
        rel = str(fp.parent.relative_to(root)) if fp.parent != root else "."
        return (fp.name, fp.suffix.lower(), rel, "—")


# ══════════════════════════════════════════════════════════════════
# Video Sorter panel
# ══════════════════════════════════════════════════════════════════

class VideoSorterPanel(_SorterPanel):

    def __init__(self, parent):
        super().__init__(parent)
        self._categories  = VIDEO_CATEGORIES
        self._cat_colors  = VIDEO_CAT_COLORS
        self._cat_icons   = VIDEO_CAT_ICONS
        self._extensions  = VIDEO_EXTENSIONS
        self.build(subtitle="screen recordings · 16:9 · 9:16 · everything else")

    def _file_type_label(self):
        return "video"

    def _classify(self, filepath: Path) -> str:
        return classify_video(filepath)

    def _tree_columns(self):
        return [
            ("name",       "File Name",  280, "w"),
            ("folder",     "Folder",     320, "w"),
            ("resolution", "Resolution", 110, "center"),
            ("ratio",      "Ratio",       90, "center"),
        ]

    def _scan_row(self, fp: Path, root: Path) -> tuple:
        dims = get_video_dimensions(str(fp))
        if dims:
            w, h = dims
            res_str   = f"{w}×{h}"
            gcd       = math.gcd(w, h)
            ratio_str = f"{w//gcd}:{h//gcd}"
        else:
            res_str   = "unknown"
            ratio_str = "—"
        rel = str(fp.parent.relative_to(root)) if fp.parent != root else "."
        return (fp.name, rel, res_str, ratio_str)

    def _error_row(self, fp: Path, root: Path) -> tuple:
        rel = str(fp.parent.relative_to(root)) if fp.parent != root else "."
        return (fp.name, rel, "error", "—")


# ══════════════════════════════════════════════════════════════════
# Document Sorter panel
# ══════════════════════════════════════════════════════════════════

class DocumentSorterPanel(_SorterPanel):

    def __init__(self, parent):
        super().__init__(parent)
        self._categories = DOCUMENT_CATEGORIES
        self._cat_colors = DOCUMENT_CAT_COLORS
        self._cat_icons  = DOCUMENT_CAT_ICONS
        self._extensions = DOCUMENT_EXTENSIONS
        self.build(subtitle="word · spreadsheets · presentations · pdf · text · code · archives")

    def _file_type_label(self):
        return "document"

    def _classify(self, filepath: Path) -> str:
        return classify_document(filepath)

    def _tree_columns(self):
        return [
            ("name",   "File Name", 320, "w"),
            ("ext",    "Ext",        70, "center"),
            ("folder", "Folder",    300, "w"),
            ("size",   "Size",       80, "e"),
        ]

    def _scan_row(self, fp: Path, root: Path) -> tuple:
        size_str = _human_size(fp.stat().st_size)
        rel      = str(fp.parent.relative_to(root)) if fp.parent != root else "."
        return (fp.name, fp.suffix.lower(), rel, size_str)

    def _error_row(self, fp: Path, root: Path) -> tuple:
        rel = str(fp.parent.relative_to(root)) if fp.parent != root else "."
        return (fp.name, fp.suffix.lower(), rel, "—")


# ══════════════════════════════════════════════════════════════════
# Audio Sorter panel
# ══════════════════════════════════════════════════════════════════

class AudioSorterPanel(_SorterPanel):

    def __init__(self, parent):
        super().__init__(parent)
        self._categories = AUDIO_CATEGORIES
        self._cat_colors = AUDIO_CAT_COLORS
        self._cat_icons  = AUDIO_CAT_ICONS
        self._extensions = AUDIO_EXTENSIONS
        self.build(subtitle="lossless · compressed · playlists · everything else")

    def _file_type_label(self):
        return "audio"

    def _classify(self, filepath: Path) -> str:
        return classify_audio(filepath)

    def _tree_columns(self):
        return [
            ("name",   "File Name", 340, "w"),
            ("ext",    "Ext",        70, "center"),
            ("folder", "Folder",    320, "w"),
            ("size",   "Size",       80, "e"),
        ]

    def _scan_row(self, fp: Path, root: Path) -> tuple:
        size_str = _human_size(fp.stat().st_size)
        rel      = str(fp.parent.relative_to(root)) if fp.parent != root else "."
        return (fp.name, fp.suffix.lower(), rel, size_str)

    def _error_row(self, fp: Path, root: Path) -> tuple:
        rel = str(fp.parent.relative_to(root)) if fp.parent != root else "."
        return (fp.name, fp.suffix.lower(), rel, "—")


# ══════════════════════════════════════════════════════════════════
# AI RENAME — text extraction helpers
# ══════════════════════════════════════════════════════════════════

AI_RENAME_EXTENSIONS = {".txt", ".pdf", ".docx", ".pptx"}

# Max chars extracted per page sent to the model — keeps context short
_PAGE_CHARS = 1200
# Max pages to summarise per file — avoid overloading context
_MAX_PAGES  = 30


def _extract_pages_txt(path: Path) -> list[str]:
    """Plain text: treat every 3 000 chars as a logical 'page'."""
    text = path.read_text(encoding="utf-8", errors="replace")
    chunk = 3000
    return [text[i:i+chunk] for i in range(0, len(text), chunk)] or [""]


def _extract_pages_pdf(path: Path) -> list[str]:
    if not _PYPDF_AVAILABLE:
        raise RuntimeError("pypdf not installed. Run: pip install pypdf")
    pages = []
    reader = _pypdf.PdfReader(str(path))
    for page in reader.pages:
        pages.append(page.extract_text() or "")
    return pages or [""]


def _extract_pages_docx(path: Path) -> list[str]:
    if not _DOCX_AVAILABLE:
        raise RuntimeError("python-docx not installed. Run: pip install python-docx")
    doc = _DocxDocument(str(path))
    # Group paragraphs into logical pages (~40 paragraphs each)
    paras = [p.text for p in doc.paragraphs if p.text.strip()]
    chunk = 40
    return ["\n".join(paras[i:i+chunk]) for i in range(0, len(paras), chunk)] or [""]


def _extract_pages_pptx(path: Path) -> list[str]:
    if not _PPTX_AVAILABLE:
        raise RuntimeError("python-pptx not installed. Run: pip install python-pptx")
    prs = _PptxPresentation(str(path))
    pages = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        pages.append("\n".join(texts))
    return pages or [""]


def extract_pages(path: Path) -> list[str]:
    """Dispatch to the right extractor based on extension."""
    ext = path.suffix.lower()
    if ext == ".txt" or ext in (".md", ".rst", ".log"):
        return _extract_pages_txt(path)
    elif ext == ".pdf":
        return _extract_pages_pdf(path)
    elif ext == ".docx":
        return _extract_pages_docx(path)
    elif ext == ".pptx":
        return _extract_pages_pptx(path)
    else:
        raise ValueError(f"Unsupported extension for AI rename: {ext}")


def _truncate(text: str, max_chars: int) -> str:
    if len(text) <= max_chars:
        return text.strip()
    return text[:max_chars].strip() + "…"


# ══════════════════════════════════════════════════════════════════
# AI RENAME — llama-server client
# ══════════════════════════════════════════════════════════════════

def _llama_complete(endpoint: str, prompt: str, max_tokens: int = 80) -> str:
    """
    POST to llama-server /completion endpoint.
    endpoint example: "http://127.0.0.1:8000"
    """
    url     = endpoint.rstrip("/") + "/completion"
    payload = json.dumps({
        "prompt":      prompt,
        "n_predict":   max_tokens,
        "temperature": 0.2,
        "top_p":       0.9,
        "stop":        ["\n", "```", "<|", "###"],
    }).encode()

    req = urllib.request.Request(
        url, data=payload,
        headers={"Content-Type": "application/json"},
        method="POST",
    )
    with urllib.request.urlopen(req, timeout=120) as resp:
        data = json.loads(resp.read())
    return data.get("content", "").strip()


def _build_rename_prompt(pages: list[str], original_name: str) -> str:
    """
    Build the prompt that asks the model for a short filename suggestion.
    Page texts are truncated and joined as a numbered list.
    """
    page_summaries = []
    for i, text in enumerate(pages[:_MAX_PAGES], 1):
        snippet = _truncate(text, _PAGE_CHARS)
        if snippet:
            page_summaries.append(f"[Page {i}]\n{snippet}")

    joined = "\n\n".join(page_summaries)

    prompt = (
        "You are a file-naming assistant. "
        "Read the document content below and suggest ONE short, descriptive filename "
        "(without extension, without quotes, no special characters except hyphens and underscores, "
        "max 6 words, use the same language as the document). "
        "Reply with ONLY the filename, nothing else.\n\n"
        f"Original filename: {original_name}\n\n"
        f"Document content:\n{joined}\n\n"
        "Suggested filename:"
    )
    return prompt


def _sanitize_suggested_name(raw: str) -> str:
    """Strip forbidden filename chars, collapse spaces to underscores."""
    # Remove characters Windows/Linux forbid in filenames
    for ch in r'\/:*?"<>|':
        raw = raw.replace(ch, "")
    raw = raw.strip().strip("'\".")
    # Collapse multiple spaces/underscores
    import re
    raw = re.sub(r"[\s_]+", "_", raw)
    raw = re.sub(r"-+", "-", raw)
    return raw[:80]  # hard cap


def ai_rename_file(filepath: Path, endpoint: str) -> str:
    """
    Full pipeline: extract pages → build prompt → call model → return clean name.
    Raises on any failure so the caller can surface the error.
    """
    pages       = extract_pages(filepath)
    prompt      = _build_rename_prompt(pages, filepath.name)
    raw         = _llama_complete(endpoint, prompt)
    suggestion  = _sanitize_suggested_name(raw)
    if not suggestion:
        raise ValueError("Model returned an empty suggestion.")
    return suggestion


# ══════════════════════════════════════════════════════════════════
# AI RENAME — panel
# ══════════════════════════════════════════════════════════════════

# Row states
_ST_PENDING  = "pending"
_ST_RUNNING  = "running"
_ST_DONE     = "done"
_ST_ERROR    = "error"
_ST_SKIPPED  = "skipped"
_ST_RENAMED  = "renamed"

_STATE_COLORS = {
    _ST_PENDING:  PALETTE["ai_pend"],
    _ST_RUNNING:  PALETTE["ai"],
    _ST_DONE:     PALETTE["ai_ok"],
    _ST_ERROR:    PALETTE["ai_err"],
    _ST_SKIPPED:  PALETTE["ai_skip"],
    _ST_RENAMED:  PALETTE["ai_ok"],
}
_STATE_ICONS = {
    _ST_PENDING:  "⏳",
    _ST_RUNNING:  "🤖",
    _ST_DONE:     "✅",
    _ST_ERROR:    "❌",
    _ST_SKIPPED:  "⏭",
    _ST_RENAMED:  "✔",
}


class AiRenamePanel(tk.Frame):
    """
    Standalone tab for AI-powered filename suggestions.

    Workflow:
      1. User picks a folder (or individual files via Browse).
      2. App scans for .txt / .pdf / .docx / .pptx files.
      3. User clicks "Analyse" — each file is sent through the
         extract → prompt → llama-server pipeline in a background thread.
      4. Results appear in a table: original name | suggested name | status.
      5. User can edit any suggestion inline, check/uncheck rows,
         then click "Apply Renames" to do the actual os.rename calls.
    """

    _COL_CHECK  = "✔"
    _COL_ORIG   = "Original filename"
    _COL_SUGG   = "Suggested name"
    _COL_STATUS = "Status"

    def __init__(self, parent):
        super().__init__(parent, bg=PALETTE["bg"])
        self._rows: list[dict] = []   # {path, suggestion, state, var(BooleanVar), iid}
        self._running  = False
        self._stop_evt = threading.Event()
        self._status   = tk.StringVar(value="Pick a folder or files to begin.")
        self._progress = tk.DoubleVar(value=0.0)
        self._endpoint = tk.StringVar(value="http://127.0.0.1:8000")
        self._build_ui()

    # ── UI construction ───────────────────────

    def _build_ui(self):
        # ── Header strip ─────────────────────
        hdr = tk.Frame(self, bg=PALETTE["panel"])
        hdr.pack(fill="x", padx=20, pady=(14, 0), ipady=10)

        tk.Label(hdr, text="🤖  AI Rename",
                 font=("Segoe UI", 13, "bold"),
                 fg=PALETTE["ai"], bg=PALETTE["panel"]).pack(side="left", padx=(14, 6))
        tk.Label(hdr, text="· reads each page · suggests a filename · you decide",
                 font=("Segoe UI", 9),
                 fg=PALETTE["subtext"], bg=PALETTE["panel"]).pack(side="left", pady=(3, 0))

        # ── Model endpoint row ────────────────
        ep_row = tk.Frame(self, bg=PALETTE["bg"])
        ep_row.pack(fill="x", padx=20, pady=(10, 0))

        tk.Label(ep_row, text="llama-server endpoint:",
                 font=("Segoe UI", 9, "bold"),
                 fg=PALETTE["subtext"], bg=PALETTE["bg"]).pack(side="left", padx=(0, 8))

        tk.Entry(ep_row, textvariable=self._endpoint,
                 font=("Segoe UI", 9), fg=PALETTE["text"],
                 bg=PALETTE["card"], bd=0, insertbackground=PALETTE["text"],
                 relief="flat", width=36,
                 ).pack(side="left", ipady=4, padx=(0, 8))

        self._ping_lbl = tk.Label(ep_row, text="",
                                  font=("Segoe UI", 9),
                                  fg=PALETTE["subtext"], bg=PALETTE["bg"])
        self._ping_lbl.pack(side="left")

        _make_btn(ep_row, "Test connection", self._test_connection,
                  PALETTE["ai"]).pack(side="left", padx=(6, 0))

        # ── Folder / file picker row ──────────
        picker = tk.Frame(self, bg=PALETTE["card"])
        picker.pack(fill="x", padx=20, pady=(10, 0), ipady=10)

        tk.Label(picker, text="Source:",
                 font=("Segoe UI", 10, "bold"),
                 fg=PALETTE["subtext"], bg=PALETTE["card"]).pack(side="left", padx=(14, 6))

        self._folder = tk.StringVar()
        tk.Entry(picker, textvariable=self._folder,
                 font=("Segoe UI", 10), fg=PALETTE["text"],
                 bg=PALETTE["panel"], bd=0, insertbackground=PALETTE["text"],
                 relief="flat", width=52,
                 ).pack(side="left", ipady=5, padx=(0, 8), fill="x", expand=True)

        _make_btn(picker, "Browse folder",
                  self._browse_folder, PALETTE["accent"]).pack(side="left", padx=(0, 4))
        _make_btn(picker, "Browse files",
                  self._browse_files,  PALETTE["accent"]).pack(side="left", padx=(0, 6))

        self._recursive = tk.BooleanVar(value=True)
        tk.Checkbutton(picker, text="Recursive",
                       variable=self._recursive,
                       font=("Segoe UI", 9), fg=PALETTE["subtext"],
                       bg=PALETTE["card"], activebackground=PALETTE["card"],
                       activeforeground=PALETTE["text"],
                       selectcolor=PALETTE["panel"],
                       ).pack(side="left", padx=(0, 14))

        # ── Progress bar ─────────────────────
        prog_row = tk.Frame(self, bg=PALETTE["bg"])
        prog_row.pack(fill="x", padx=20, pady=(8, 0))

        ttk.Progressbar(prog_row, variable=self._progress,
                        maximum=100, mode="determinate",
                        ).pack(side="left", fill="x", expand=True)
        tk.Label(prog_row, textvariable=self._status,
                 font=("Segoe UI", 9), fg=PALETTE["subtext"],
                 bg=PALETTE["bg"], anchor="w",
                 ).pack(side="left", padx=(10, 0))

        # ── Results treeview ─────────────────
        tree_frame = tk.Frame(self, bg=PALETTE["bg"])
        tree_frame.pack(fill="both", expand=True, padx=20, pady=(12, 0))
        tree_frame.rowconfigure(0, weight=1)
        tree_frame.columnconfigure(0, weight=1)

        cols = ("check", "orig", "sugg", "status")
        self._tree = ttk.Treeview(tree_frame, columns=cols,
                                  show="headings", selectmode="browse")
        self._tree.heading("check",  text="✔",           anchor="center")
        self._tree.heading("orig",   text="Original",    anchor="w")
        self._tree.heading("sugg",   text="Suggested name (double-click to edit)", anchor="w")
        self._tree.heading("status", text="Status",      anchor="center")

        self._tree.column("check",  width=36,  minwidth=36,  stretch=False, anchor="center")
        self._tree.column("orig",   width=280, minwidth=160, stretch=True,  anchor="w")
        self._tree.column("sugg",   width=280, minwidth=160, stretch=True,  anchor="w")
        self._tree.column("status", width=110, minwidth=90,  stretch=False, anchor="center")

        vsb = ttk.Scrollbar(tree_frame, orient="vertical",   command=self._tree.yview)
        hsb = ttk.Scrollbar(tree_frame, orient="horizontal", command=self._tree.xview)
        self._tree.configure(yscrollcommand=vsb.set, xscrollcommand=hsb.set)

        vsb.grid(row=0, column=1, sticky="ns")
        hsb.grid(row=1, column=0, sticky="ew")
        self._tree.grid(row=0, column=0, sticky="nsew")

        # Double-click → inline edit of suggested name
        self._tree.bind("<Double-1>", self._on_double_click)
        # Single-click on check column → toggle include
        self._tree.bind("<ButtonRelease-1>", self._on_click)

        # ── Action bar ───────────────────────
        bar = tk.Frame(self, bg=PALETTE["panel"])
        bar.pack(fill="x", padx=20, pady=(10, 16))

        tk.Label(bar, text="Actions:",
                 font=("Segoe UI", 9, "bold"),
                 fg=PALETTE["subtext"], bg=PALETTE["panel"]).pack(side="left", padx=(14, 10))

        self._scan_btn = _make_btn(bar, "▶  Scan files",
                                   self._start_scan, PALETTE["accent3"])
        self._scan_btn.pack(side="left", padx=(0, 6))

        self._stop_btn = _make_btn(bar, "■  Stop",
                                   self._stop_scan, PALETTE["accent2"])
        self._stop_btn.pack(side="left", padx=(0, 6))
        self._stop_btn.configure(state="disabled")

        self._analyse_btn = _make_btn(bar, "🤖  Analyse all",
                                      self._start_analyse, PALETTE["ai"])
        self._analyse_btn.pack(side="left", padx=(0, 6))
        self._analyse_btn.configure(state="disabled")

        self._apply_btn = _make_btn(bar, "✔  Apply renames",
                                    self._apply_renames, PALETTE["ai_ok"])
        self._apply_btn.pack(side="left", padx=(0, 6))
        self._apply_btn.configure(state="disabled")

        self._clear_btn = _make_btn(bar, "🗑  Clear",
                                    self._clear_all, PALETTE["subtext"])
        self._clear_btn.pack(side="left")

        tk.Label(bar,
                 text="Supports: .txt  .pdf  .docx  .pptx",
                 font=("Segoe UI", 8), fg=PALETTE["subtext"], bg=PALETTE["panel"],
                 ).pack(side="right", padx=14)

    # ── Connection test ───────────────────────

    def _test_connection(self):
        endpoint = self._endpoint.get().strip()
        self._ping_lbl.configure(text="testing…", fg=PALETTE["subtext"])
        self.update_idletasks()
        try:
            url = endpoint.rstrip("/") + "/health"
            req = urllib.request.Request(url, method="GET")
            with urllib.request.urlopen(req, timeout=5) as r:
                r.read()
            self._ping_lbl.configure(text="✅ connected", fg=PALETTE["ai_ok"])
        except Exception as e:
            self._ping_lbl.configure(text=f"❌ {e}", fg=PALETTE["ai_err"])

    # ── Browse ────────────────────────────────

    def _browse_folder(self):
        folder = filedialog.askdirectory(title="Select folder to scan")
        if folder:
            self._folder.set(folder)

    def _browse_files(self):
        files = filedialog.askopenfilenames(
            title="Select files",
            filetypes=[("Supported", "*.txt *.pdf *.docx *.pptx"),
                       ("All files", "*.*")])
        if files:
            # Show parent of first file in the entry, store paths directly
            self._folder.set(str(Path(files[0]).parent))
            self._load_file_list(list(files))

    # ── Scan ──────────────────────────────────

    def _start_scan(self):
        folder = self._folder.get().strip()
        if not folder or not os.path.isdir(folder):
            messagebox.showerror("No folder", "Please select a valid folder first.")
            return
        self._clear_all()
        self._stop_evt.clear()
        self._running = True
        self._scan_btn.configure(state="disabled")
        self._stop_btn.configure(state="normal")
        self._analyse_btn.configure(state="disabled")
        self._apply_btn.configure(state="disabled")
        threading.Thread(target=self._scan_worker, args=(folder,), daemon=True).start()

    def _scan_worker(self, folder: str):
        root = Path(folder)
        self.after(0, self._status.set, "Scanning for supported files…")
        try:
            if self._recursive.get():
                files = [p for p in root.rglob("*")
                         if p.is_file() and p.suffix.lower() in AI_RENAME_EXTENSIONS]
            else:
                files = [p for p in root.iterdir()
                         if p.is_file() and p.suffix.lower() in AI_RENAME_EXTENSIONS]
        except Exception as e:
            self.after(0, self._status.set, f"Scan error: {e}")
            self.after(0, self._finish_scan)
            return

        if not files:
            self.after(0, self._status.set, "No supported files found (.txt .pdf .docx .pptx).")
            self.after(0, self._finish_scan)
            return

        self.after(0, self._load_file_list, [str(f) for f in files])
        self.after(0, self._status.set, f"Found {len(files)} file(s). Ready to analyse.")
        self.after(0, self._finish_scan)

    def _load_file_list(self, file_paths: list[str]):
        self._rows.clear()
        for child in self._tree.get_children():
            self._tree.delete(child)

        for fp_str in file_paths:
            fp  = Path(fp_str)
            var = tk.BooleanVar(value=True)
            iid = self._tree.insert("", "end", values=(
                "☑", fp.name, "—", _STATE_ICONS[_ST_PENDING] + " pending"
            ))
            row = {"path": fp, "suggestion": "", "state": _ST_PENDING,
                   "var": var, "iid": iid}
            self._rows.append(row)
            self._tree.tag_configure("pending", foreground=PALETTE["ai_pend"])
            self._tree.item(iid, tags=("pending",))

        if self._rows:
            self._analyse_btn.configure(state="normal")

    def _finish_scan(self):
        self._running = False
        self._scan_btn.configure(state="normal")
        self._stop_btn.configure(state="disabled")

    def _stop_scan(self):
        self._stop_evt.set()

    # ── Analyse ───────────────────────────────

    def _start_analyse(self):
        pending = [r for r in self._rows if r["state"] in (_ST_PENDING, _ST_ERROR)]
        if not pending:
            messagebox.showinfo("Nothing to do", "No pending files to analyse.")
            return
        self._stop_evt.clear()
        self._running = True
        self._analyse_btn.configure(state="disabled")
        self._apply_btn.configure(state="disabled")
        self._scan_btn.configure(state="disabled")
        self._stop_btn.configure(state="normal")
        threading.Thread(target=self._analyse_worker,
                         args=(pending,), daemon=True).start()

    def _analyse_worker(self, rows: list[dict]):
        endpoint = self._endpoint.get().strip()
        total    = len(rows)

        for i, row in enumerate(rows):
            if self._stop_evt.is_set():
                self.after(0, self._status.set, f"Stopped at {i}/{total}.")
                break

            self.after(0, self._progress.set, (i / total) * 100)
            self.after(0, self._status.set,
                       f"[{i+1}/{total}]  Analysing {row['path'].name}…")
            self.after(0, self._set_row_state, row, _ST_RUNNING, "…")

            try:
                suggestion = ai_rename_file(row["path"], endpoint)
                row["suggestion"] = suggestion
                row["state"]      = _ST_DONE
                self.after(0, self._set_row_state, row, _ST_DONE, suggestion)
            except Exception as e:
                row["state"] = _ST_ERROR
                self.after(0, self._set_row_state, row, _ST_ERROR, str(e)[:60])

        self.after(0, self._progress.set, 100)
        done  = sum(1 for r in self._rows if r["state"] == _ST_DONE)
        errs  = sum(1 for r in self._rows if r["state"] == _ST_ERROR)
        self.after(0, self._status.set,
                   f"Analysis complete — {done} suggestions, {errs} error(s).")
        self.after(0, self._finish_analyse)

    def _finish_analyse(self):
        self._running = False
        self._scan_btn.configure(state="normal")
        self._stop_btn.configure(state="disabled")
        self._analyse_btn.configure(state="normal")
        has_done = any(r["state"] == _ST_DONE and r["var"].get() for r in self._rows)
        if has_done:
            self._apply_btn.configure(state="normal")

    def _set_row_state(self, row: dict, state: str, suggestion: str):
        color = _STATE_COLORS.get(state, PALETTE["text"])
        icon  = _STATE_ICONS.get(state, "")
        check = "☑" if row["var"].get() else "☐"
        self._tree.item(row["iid"], values=(
            check, row["path"].name, suggestion, f"{icon} {state}"
        ))
        # Re-tag for colour
        tag = f"state_{state}"
        self._tree.tag_configure(tag, foreground=color)
        self._tree.item(row["iid"], tags=(tag,))
        if suggestion and state == _ST_DONE:
            row["suggestion"] = suggestion

    # ── Inline edit (double-click on suggestion column) ───────────

    def _on_double_click(self, event):
        region = self._tree.identify_region(event.x, event.y)
        col    = self._tree.identify_column(event.x)
        iid    = self._tree.identify_row(event.y)
        if region != "cell" or col != "#3" or not iid:
            return

        row = next((r for r in self._rows if r["iid"] == iid), None)
        if row is None or row["state"] not in (_ST_DONE, _ST_ERROR):
            return

        # Get bounding box of the cell
        x, y, w, h = self._tree.bbox(iid, col)

        current = row["suggestion"]
        var = tk.StringVar(value=current)

        entry = tk.Entry(self._tree, textvariable=var,
                         font=("Segoe UI", 9), fg=PALETTE["text"],
                         bg=PALETTE["sel"], bd=0,
                         insertbackground=PALETTE["text"], relief="flat")
        entry.place(x=x, y=y, width=w, height=h)
        entry.focus_set()
        entry.select_range(0, "end")

        def _commit(evt=None):
            new_val = _sanitize_suggested_name(var.get())
            row["suggestion"] = new_val
            row["state"]      = _ST_DONE
            self._set_row_state(row, _ST_DONE, new_val)
            entry.destroy()
            self._apply_btn.configure(state="normal")

        def _cancel(evt=None):
            entry.destroy()

        entry.bind("<Return>",  _commit)
        entry.bind("<Tab>",     _commit)
        entry.bind("<Escape>",  _cancel)
        entry.bind("<FocusOut>", _commit)

    # ── Toggle include checkbox (click on check column) ───────────

    def _on_click(self, event):
        col = self._tree.identify_column(event.x)
        iid = self._tree.identify_row(event.y)
        if col != "#1" or not iid:
            return
        row = next((r for r in self._rows if r["iid"] == iid), None)
        if row is None:
            return
        row["var"].set(not row["var"].get())
        check = "☑" if row["var"].get() else "☐"
        vals = list(self._tree.item(iid, "values"))
        vals[0] = check
        self._tree.item(iid, values=vals)
        # Enable/disable apply based on whether anything is checked+done
        has_done = any(r["state"] == _ST_DONE and r["var"].get() for r in self._rows)
        self._apply_btn.configure(state="normal" if has_done else "disabled")

    # ── Apply renames ─────────────────────────

    def _apply_renames(self):
        targets = [r for r in self._rows
                   if r["state"] == _ST_DONE and r["var"].get() and r["suggestion"]]
        if not targets:
            messagebox.showinfo("Nothing to do",
                                "No checked files with suggestions to rename.")
            return

        preview = "\n".join(
            f"  {r['path'].name}  →  {r['suggestion']}{r['path'].suffix}"
            for r in targets[:20]
        )
        if len(targets) > 20:
            preview += f"\n  … and {len(targets)-20} more"

        if not messagebox.askyesno("Confirm renames",
                                   f"Rename {len(targets)} file(s)?\n\n{preview}"):
            return

        done = 0
        errors = []
        for row in targets:
            new_name = row["suggestion"] + row["path"].suffix
            new_path = row["path"].parent / new_name
            # Avoid collision
            stem, ext = row["suggestion"], row["path"].suffix
            n = 1
            while new_path.exists() and new_path.resolve() != row["path"].resolve():
                new_path = row["path"].parent / f"{stem}_{n}{ext}"
                n += 1
            try:
                row["path"].rename(new_path)
                row["path"]  = new_path
                row["state"] = _ST_RENAMED
                self._set_row_state(row, _ST_RENAMED, row["suggestion"])
                done += 1
            except Exception as e:
                errors.append(f"{row['path'].name}: {e}")
                row["state"] = _ST_ERROR
                self._set_row_state(row, _ST_ERROR, str(e)[:60])

        msg = f"✓  {done} file(s) renamed successfully."
        if errors:
            msg += f"\n\n✗  {len(errors)} error(s):\n" + "\n".join(f"  • {e}" for e in errors[:10])
        messagebox.showinfo("Done", msg)
        self._status.set(f"Renamed {done}/{len(targets)} files.")

    # ── Clear ─────────────────────────────────

    def _clear_all(self):
        self._rows.clear()
        for child in self._tree.get_children():
            self._tree.delete(child)
        self._progress.set(0)
        self._status.set("Cleared. Pick a folder or files to begin.")
        self._analyse_btn.configure(state="disabled")
        self._apply_btn.configure(state="disabled")


# ══════════════════════════════════════════════════════════════════
# Main application window
# ══════════════════════════════════════════════════════════════════

class MediaSortApp(tk.Tk):

    def __init__(self):
        super().__init__()
        self.title("MediaSort")
        self.geometry("1080x780")
        self.minsize(840, 600)
        self.configure(bg=PALETTE["bg"])

        _apply_shared_styles(self)
        self._build_ui()

    def _build_ui(self):
        # ── App header ───────────────────────────
        hdr = tk.Frame(self, bg=PALETTE["bg"])
        hdr.pack(fill="x", padx=20, pady=(18, 0))

        tk.Label(
            hdr, text="MediaSort",
            font=("Segoe UI", 22, "bold"),
            fg=PALETTE["accent"], bg=PALETTE["bg"],
        ).pack(side="left")

        tk.Label(
            hdr, text=" — images, videos, documents & audio, organised",
            font=("Segoe UI", 11),
            fg=PALETTE["subtext"], bg=PALETTE["bg"],
        ).pack(side="left", pady=(8, 0))

        # ── Mode switcher ────────────────────────
        switcher = tk.Frame(self, bg=PALETTE["bg"])
        switcher.pack(fill="x", padx=20, pady=(12, 0))

        self._mode = tk.StringVar(value="images")

        _tab_defs = [
            ("Images",    "images",    "🖼"),
            ("Videos",    "videos",    "🎬"),
            ("Documents", "documents", "📄"),
            ("Audio",     "audio",     "🎵"),
            ("AI Rename", "airename",  "🤖"),
        ]
        self._tab_btns: dict[str, tk.Button] = {}

        def _select_tab(value):
            self._mode.set(value)
            _refresh_tabs()

        for text, value, icon in _tab_defs:
            b = tk.Button(
                switcher, text=f"  {icon}  {text}  ",
                command=lambda v=value: _select_tab(v),
                font=("Segoe UI", 10, "bold"),
                fg=PALETTE["text"], bg=PALETTE["card"],
                activeforeground=PALETTE["text"], activebackground=PALETTE["hover"],
                relief="flat", bd=0, cursor="hand2", padx=8, pady=6,
            )
            b.pack(side="left", padx=(0, 4))
            self._tab_btns[value] = b

        sep = tk.Frame(self, bg=PALETTE["border"], height=1)
        sep.pack(fill="x", padx=20, pady=(10, 0))

        # ── Panels ───────────────────────────────
        self._img_panel  = ImageSorterPanel(self)
        self._vid_panel  = VideoSorterPanel(self)
        self._doc_panel  = DocumentSorterPanel(self)
        self._aud_panel  = AudioSorterPanel(self)
        self._ai_panel   = AiRenamePanel(self)

        self._panels = {
            "images":    self._img_panel,
            "videos":    self._vid_panel,
            "documents": self._doc_panel,
            "audio":     self._aud_panel,
            "airename":  self._ai_panel,
        }

        def _refresh_tabs():
            mode = self._mode.get()
            for key, panel in self._panels.items():
                panel.pack_forget()
            self._panels[mode].pack(fill="both", expand=True)
            for key, btn in self._tab_btns.items():
                if key == mode:
                    btn.configure(bg=PALETTE["accent"], fg=PALETTE["bg"])
                else:
                    btn.configure(bg=PALETTE["card"], fg=PALETTE["text"])

        _refresh_tabs()   # show images panel by default


# ══════════════════════════════════════════════════════════════════
# Entry point
# ══════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    app = MediaSortApp()
    app.mainloop()
